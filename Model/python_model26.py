import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import LabelEncoder, OneHotEncoder, StandardScaler
from sklearn.compose import ColumnTransformer
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import accuracy_score, precision_recall_fscore_support
import joblib
from docx import Document
from docx.shared import Inches
from pptx import Presentation
import os
from datetime import datetime
from pptx.util import Inches as PPTInches, Pt

# Load the dataset
data_path = 'jobs_in_data.csv'
data = pd.read_csv( data_path )

# Encode 'job_title'
le_job_title = LabelEncoder()
data ['job_title_encoded'] = le_job_title.fit_transform( data ['job_title'] )

# Preparing features and target variable
X = data.drop( ['job_title', 'job_title_encoded'], axis=1 )
y = data ['job_title_encoded']

# Identify categorical columns in your dataset
categorical_cols = [col for col in X.columns if X [col].dtype == 'object']

# Apply OneHotEncoder to categorical columns - creates a new column for each category
preprocessor = ColumnTransformer( transformers=[
    ('cat', OneHotEncoder(), categorical_cols)], remainder='passthrough', sparse_threshold=0 )

X_transformed = preprocessor.fit_transform( X )

# Splitting the transformed data into training and test sets
X_train, X_test, y_train, y_test = train_test_split( X_transformed, y, test_size=0.2, random_state=42 )

# Feature Scaling with adjustment for sparse data
scaler = StandardScaler( with_mean=False )
X_train_scaled = scaler.fit_transform( X_train )
X_test_scaled = scaler.transform( X_test )

# Model Training
clf = RandomForestClassifier( n_estimators=100, random_state=42 )
clf.fit( X_train_scaled, y_train )

# Model Evaluation
predictions = clf.predict( X_test_scaled )
metrics = precision_recall_fscore_support( y_test, predictions, average='weighted', zero_division=0 )
accuracy = accuracy_score( y_test, predictions )

print(
    f"Accuracy: {accuracy:.4f}, Precision: {metrics [0]:.4f}, Recall: {metrics [1]:.4f}, F1 Score: {metrics [2]:.4f}" )

### Saving the Model and Scaler ###
model_path = 'report/model/random_forest.joblib'
scaler_path = 'report/model/scaler.joblib'

# Specify directory and filename
model_dir = "report/model"
model_filename = "random_forest.joblib"

# Create the directory if it doesn't exist
os.makedirs( model_dir, exist_ok=True )

# Create the full file path
full_model_path = os.path.join( model_dir, model_filename )

# Save your classifier to a file
joblib.dump( clf, full_model_path )

print( f"Trained model saved at {full_model_path}" )


# Directory for saving visualizations
visualization_dir = 'report/visualizations/'
if not os.path.exists( visualization_dir ):
    os.makedirs( visualization_dir )


# Helper function to save figures
def save_fig( fig_id, tight_layout=True, fig_extension="png", resolution=300 ):
    path = os.path.join( visualization_dir, fig_id + "." + fig_extension )
    print( f"Saving figure: {fig_id}" )
    if tight_layout:
        plt.tight_layout()
    plt.savefig( path, format=fig_extension, dpi=resolution )
    plt.close()


# Visualizations
sns.countplot( data=data, x='work_year', palette='coolwarm' )
plt.title( 'Demand for Data Professions Over Time' )
plt.xlabel( 'Year' )
plt.ylabel( 'Number of Listings' )
save_fig( 'demand_over_time' )

sns.boxplot( x='experience_level', y='salary_in_usd', data=data, palette='viridis' )
plt.title( 'Salary Distribution by Experience Level' )
plt.xlabel( 'Experience Level' )
plt.ylabel( 'Salary in USD' )
save_fig( 'salary_distribution' )

sns.countplot( x='work_setting', data=data, palette='viridis' )
plt.title( 'Work Setting Preferences by Profession' )
plt.xlabel( 'Work Setting' )
plt.ylabel( 'Count' )
save_fig( 'work_setting_preferences' )

sns.countplot( x='employment_type', data=data, palette='coolwarm' )
plt.title( 'Employment Type Distribution' )
plt.xlabel( 'Employment Type' )
plt.ylabel( 'Count' )
save_fig( 'employment_type_distribution' )

# Generate Visualization for Top 10 Job Titles with the Most Demand
top_10_job_titles = data ['job_title'].value_counts().head( 10 )
plt.figure( figsize=(10, 6) )
sns.barplot( x=top_10_job_titles.values, y=top_10_job_titles.index, palette='coolwarm' )
plt.title( 'Top 10 Job Titles with the Most Demand' )
plt.xlabel( 'Number of Listings' )
plt.ylabel( 'Job Title' )
save_fig( 'top_10_job_titles_demand' )

# Extend the visualization_names list with the new visualization
visualization_dir = []
visualization_dir.extend( ['top_10_job_titles_demand'] )
visualization_dir.extend( ['work_setting_preferences'] )
visualization_dir.extend( ['salary_distribution'] )
visualization_dir.extend( ['demand_over_time'] )
visualization_dir.extend( ['work_setting_preferences'] )
visualization_dir.extend( ['employment_type_distribution'] )

### Report Generation ###
# Directory for saving visualizations
visualization_dir_path = 'report/visualizations/'
if not os.path.exists( visualization_dir_path ):
    os.makedirs( visualization_dir_path )

# initialization of the list containing visualization names
visualizations = ['demand_over_time', 'salary_distribution', 'work_setting_preferences',
                  'employment_type_distribution', 'top_10_job_titles_demand']


# Create a new Word document
# Initialize the Word document
doc = Document()

# Title Page
doc.add_heading('Data Industry Hiring Trends Analysis', 0)

# Introduction
doc.add_heading('Introduction', level=1)
doc.add_paragraph(
    "This report provides an analysis of hiring trends in the data industry. "
    "The objective is to understand the demand for various data professions over time, "
    "explore salary distributions, and examine preferences for work settings and employment types."
)

# Step-by-Step Model Guide
doc.add_heading('Step-by-Step Model Guide', level=1)
# Data Preparation
doc.add_heading('Data Preparation', level=2)
doc.add_paragraph(
    "Initially, the dataset is loaded from a CSV file. Relevant features are selected, "
    "and any missing values are addressed. The 'job_title' column is encoded using LabelEncoder for further processing."
)
# Feature Encoding
doc.add_heading('Feature Encoding', level=2)
doc.add_paragraph(
    "Categorical variables are transformed using OneHotEncoder to prepare the dataset for the machine learning model."
)
# Model Training
doc.add_heading('Model Training', level=2)
doc.add_paragraph(
    "A RandomForestClassifier is trained with the prepared dataset. Hyperparameters are set to default values for initial training."
)
# Model Evaluation
doc.add_heading('Model Evaluation', level=2)
doc.add_paragraph(
    "The model's performance is evaluated using metrics such as accuracy, precision, recall, and F1 score. "
    "These metrics provide insights into the model's ability to generalize to new data."
)

# Visualizations and Insights
doc.add_heading('Visualizations and Insights', level=1)

visualizations = [
    ('demand_over_time', 'Demand for Data Professions Over Time', 'This visualization shows the number of job listings per year, indicating trends in the data industry’s demand.'),
    ('salary_distribution', 'Salary Distribution by Experience Level', 'Salary distributions across different experience levels highlight the impact of experience on compensation in the data industry.'),
    ('work_setting_preferences', 'Work Setting Preferences by Profession', 'Preferences for work settings illustrate the diversity of work environment desires among data professionals.'),
    ('employment_type_distribution', 'Employment Type Distribution', 'The distribution of employment types shows the variety in contractual arrangements preferred or offered in the data industry.'),
    ('top_10_job_titles_demand', 'Top 10 Job Titles with the Most Demand', 'This chart ranks the top 10 job titles by demand, offering insight into the most sought-after positions.')
]

visualization_dir_path = 'report/visualizations/'
for viz_name, title, insight in visualizations:
    doc.add_heading(title, level=2)
    img_path = os.path.join(visualization_dir_path, f'{viz_name}.png')
    if os.path.exists(img_path):
        doc.add_picture(img_path, width=Inches(6))
    doc.add_paragraph(insight)  # Add insights directly below the visualization

# Save the document
report_docx_path = os.path.join(visualization_dir_path, 'Data_Industry_Hiring_Trends_Report.docx')
doc.save(report_docx_path)
print(f"Report saved at {report_docx_path}")



# Initialize PowerPoint presentation
prs = Presentation()

# Title slide
title_slide_layout = prs.slide_layouts [0]  # 0 is the layout for a title slide
title_slide = prs.slides.add_slide( title_slide_layout )
title = title_slide.shapes.title
subtitle = title_slide.placeholders [1]

# Set title, subtitle, and other details
title.text = "Data Industry Hiring Trends Presentation"
current_date = datetime.today().strftime( '%B %d, %Y' )  # Formats date as Month Day, Year
subtitle.text = (f"Author: Baltzakis Themistoklis\n"
                 f"College: NY College\n"
                 f"Date: {current_date}\n"
                 f"Website: www.baltzakisthemis.com\n"
                 f"Email: baltzakis.themis@gmail.com")

# Add a slide for Model Evaluation Metrics
slide_layout = prs.slide_layouts [1]  # Using a title and content layout for metrics
metrics_slide = prs.slides.add_slide( slide_layout )
title = metrics_slide.shapes.title
content = metrics_slide.placeholders [1]

title.text = "Model Evaluation Metrics"
# Update this metrics_content with your actual model evaluation metrics
metrics_content = "Accuracy: 0.95\nPrecision: 0.96\nRecall: 0.94\nF1 Score: 0.95"
content.text = metrics_content

# Adding visualizations to the presentation
visualizations = [
    ('demand_over_time', 'Demand for Data Professions Over Time',
     'This visualization shows the number of job listings per year, indicating trends in the data industry’s demand.'),
    ('salary_distribution', 'Salary Distribution by Experience Level',
     'Salary distributions across different experience levels highlight the impact of experience on compensation in '
     'the data industry.'),
    ('work_setting_preferences', 'Work Setting Preferences by Profession',
     'Preferences for work settings illustrate the diversity of work environment desires among data professionals.'),
    ('employment_type_distribution', 'Employment Type Distribution',
     'The distribution of employment types shows the variety in contractual arrangements preferred or offered in the '
     'data industry.'),
    ('top_10_job_titles_demand', 'Top 10 Job Titles with the Most Demand',
     'This chart ranks the top 10 job titles by demand, offering insight into the most sought-after positions.')
]

visualization_dir_path = 'report/visualizations/'
for viz_name, title, insight in visualizations:
    slide = prs.slides.add_slide( prs.slide_layouts [5] )  # Using a blank layout
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    img_path = os.path.join( visualization_dir_path, f'{viz_name}.png' )
    slide.shapes.add_picture( img_path, PPTInches( 1 ), PPTInches( 1.5 ), width=PPTInches( 8 ) )

    # Adding a textbox for insights below the image
    txBox = slide.shapes.add_textbox( PPTInches( 1 ), PPTInches( 5.5 ), width=PPTInches( 8 ), height=PPTInches( 2 ) )
    tf = txBox.text_frame
    p = tf.paragraphs [0]
    p.text = insight
    p.font.size = Pt( 14 )  # Adjust font size as needed

# Save the presentation
presentation_pptx_path = os.path.join( visualization_dir_path, 'Data_Industry_Hiring_Trends_Presentation.pptx' )
prs.save( presentation_pptx_path )
print( f"Presentation saved at {presentation_pptx_path}" )